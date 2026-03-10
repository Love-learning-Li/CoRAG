#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
-------------------------------------------------------------------------
This file is part of the RAGSDK project.
Copyright (c) 2025 Huawei Technologies Co.,Ltd.

RAGSDK is licensed under Mulan PSL v2.
You can use this software according to the terms and conditions of the Mulan PSL v2.
You may obtain a copy of Mulan PSL v2 at:

         http://license.coscl.org.cn/MulanPSL2

THIS SOFTWARE IS PROVIDED ON AN "AS IS" BASIS, WITHOUT WARRANTIES OF ANY KIND,
EITHER EXPRESS OR IMPLIED, INCLUDING BUT NOT LIMITED TO NON-INFRINGEMENT,
MERCHANTABILITY OR FIT FOR A PARTICULAR PURPOSE.
See the Mulan PSL v2 for more details.
-------------------------------------------------------------------------
"""

import itertools

from loguru import logger
from paddleocr import PaddleOCR
from pptx import Presentation
from langchain_core.documents import Document
from langchain_community.document_loaders.base import BaseLoader
from tqdm import tqdm

from mx_rag.document.loader.base_loader import BaseLoader as mxBaseLoader
from mx_rag.llm import Img2TextLLM
from mx_rag.utils.common import validate_params, Lang, BOOL_TYPE_CHECK_TIP
from mx_rag.utils.file_check import SecFileCheck


class PowerPointLoader(BaseLoader, mxBaseLoader):
    EXTENSION = (".pptx",)
    MAX_SIZE = 100 * 1024 * 1024
    MAX_TABLE_ROW = 100
    MAX_TABLE_COL = 50

    @validate_params(
        vlm=dict(validator=lambda x: isinstance(x, Img2TextLLM) or x is None,
                 message="param must be instance of Img2TextLLM or None"),
        lang=dict(validator=lambda x: isinstance(x, Lang), message="param must be instance of Lang"),
        enable_ocr=dict(validator=lambda x: isinstance(x, bool), message=BOOL_TYPE_CHECK_TIP)
    )
    def __init__(self, file_path, vlm: Img2TextLLM = None, lang=Lang.CH, enable_ocr=False):
        super().__init__(file_path)
        self.enable_ocr = enable_ocr
        self.vlm = vlm

        try:
            self.ocr = PaddleOCR(use_angle_cls=True, lang=lang.value, show_log=False) if enable_ocr else None
        except ImportError as err:
            raise ImportError(f"init ocr failed due to import error, {err}") from err
        except MemoryError as err:
            raise MemoryError(f"init ocr failed due to memory error, {err}") from err
        except Exception as err:
            raise Exception(f"init ocr failed, {err}") from err

    def lazy_load(self):
        self._check_file_valid()
        try:
            return self._load_ppt()
        except ValueError as ve:
            logger.error(f"Value error: {ve}")
            return iter([])
        except Exception as err:
            logger.error(f"load '{self.file_path}' failed, {err}")
            return iter([])

    def _check_file_valid(self):
        SecFileCheck(self.file_path, self.MAX_SIZE).check()
        if not self.file_path.endswith(self.EXTENSION):
            raise TypeError("file type not correct")
        if self._is_zip_bomb():
            raise ValueError(f"'{self.file_path}' is a risk of zip bombs")

    def _load_merged_cell(self, data, cell, row, col):
        span_height = cell.span_height
        span_width = cell.span_width
        for span_row in range(row, row + span_height):
            for span_col in range(col, col + span_width):
                data[span_row][span_col] = cell.text

    def _load_table(self, table):
        if (len(table.rows) > self.MAX_TABLE_ROW) or (len(table.columns) > self.MAX_TABLE_COL):
            logger.warning(f"can not load table over {self.MAX_TABLE_ROW} rows or {self.MAX_TABLE_COL} cols")
        rows = min(len(table.rows), self.MAX_TABLE_ROW)
        cols = min(len(table.columns), self.MAX_TABLE_COL)
        # 初始化一个二维列表来存储表格数据
        data = [["" for _ in range(cols)] for _ in range(rows)]
        for row in range(rows):
            for col in range(cols):
                cell = table.cell(row, col)
                if not cell.text:
                    continue
                if not cell.is_merge_origin:
                    data[row][col] = cell.text
                    continue
                self._load_merged_cell(data, cell, row, col)

        return itertools.chain.from_iterable(data)

    def _load_image_text(self, image_bytes):
        if not self._verify_image_size(image_bytes):
            return None
        result = self.ocr.ocr(image_bytes, cls=True)
        try:
            res = [line[1][0] for line in result[0]]
            return res
        except TypeError as err:
            logger.info(f"can not load text from image, {err}")
            return None

    def _load_slide(self, slide):
        slide_text, img_base64_list, image_summaries = [], [], []
        for shape in slide.shapes:
            if hasattr(shape, "image") and self.enable_ocr and self.ocr is not None:
                image_data = shape.image.blob
                img_text = self._load_image_text(image_data)
                if img_text is not None:
                    slide_text.extend(img_text)

            if hasattr(shape, "image") and self.vlm:
                image_data = shape.image.blob
                img_base64, image_summary = self._interpret_image(image_data, self.vlm)
                img_base64_list.extend([img_base64] if image_summary and img_base64 else [])
                image_summaries.extend([image_summary] if image_summary and img_base64 else [])

            if shape.has_table:
                table = shape.table
                table_text = self._load_table(table)
                slide_text.extend(table_text)

            if shape.has_text_frame:
                slide_text.append(shape.text_frame.text.replace("\n", " "))
        return " ".join(slide_text), img_base64_list, image_summaries

    def _load_ppt(self):
        prs = Presentation(self.file_path)
        total_slides = len(prs.slides)

        # Only show progress bar if there are more than 5 slides
        for slide in tqdm(prs.slides, desc="Processing slides", total=total_slides, disable=total_slides < 5):
            slide_text, img_base64_list, image_summaries = self._load_slide(slide)
            for img_base64, image_summary in zip(img_base64_list, image_summaries):
                yield Document(page_content=image_summary, metadata={"source": self.file_path,
                                                                     "image_base64": img_base64, "type": "image"})

            yield Document(page_content=slide_text, metadata={"source": self.file_path, "type": "text"})
